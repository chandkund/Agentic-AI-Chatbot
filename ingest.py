"""
Robust Document Ingestion Service with ChromaDB Integration
"""

import os
import uuid
import shutil
from pathlib import Path
from typing import List, Optional, Dict, Any
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Configuration
UPLOAD_DIR = Path("uploads")
DB_DIR = Path("chroma_db")
COLLECTION_NAME = "documents"
EMBED_MODEL_NAME = "all-MiniLM-L6-v2"
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200

# Ensure directories exist
UPLOAD_DIR.mkdir(exist_ok=True, parents=True)
DB_DIR.mkdir(exist_ok=True, parents=True)

class DocumentIngester:
    def __init__(self):
        """Initialize the document ingester with ChromaDB backend"""
        # Initialize embedding model
        self.embedding_fn = self._initialize_embeddings()
        
        # Initialize ChromaDB components
        self.client = self._initialize_chroma_client()
        self.collection = self._initialize_collection()
        
        # Initialize text splitter
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP
        )

    def _initialize_embeddings(self):
        """Initialize sentence transformer embeddings"""
        return embedding_functions.SentenceTransformerEmbeddingFunction(
            model_name=EMBED_MODEL_NAME
        )

    def _get_chroma_settings(self) -> Settings:
        """Return consistent ChromaDB settings"""
        return Settings(
            anonymized_telemetry=False,
            allow_reset=True,
            is_persistent=True
        )

    def _clean_chroma_db(self):
        """Completely clean up ChromaDB directory and reset state"""
        try:
            # Clear any existing Chroma state
            if hasattr(chromadb, '_instance'):
                del chromadb._instance
            
            # Clear the directory
            shutil.rmtree(DB_DIR, ignore_errors=True)
            DB_DIR.mkdir(exist_ok=True, parents=True)
        except Exception as e:
            print(f"Error cleaning ChromaDB: {e}")
            raise

    def _initialize_chroma_client(self) -> chromadb.PersistentClient:
        """Initialize ChromaDB client with proper error handling"""
        max_retries = 3
        last_exception = None
        
        for attempt in range(max_retries):
            try:
                # Clear any existing client instances
                if hasattr(chromadb, '_instance'):
                    del chromadb._instance
                    
                return chromadb.PersistentClient(
                    path=str(DB_DIR),
                    settings=self._get_chroma_settings()
                )
            except Exception as e:
                last_exception = e
                print(f"Chroma client initialization attempt {attempt + 1} failed: {e}")
                if attempt < max_retries - 1:
                    self._clean_chroma_db()
        
        print("All Chroma client initialization attempts failed")
        raise RuntimeError(f"Failed to initialize Chroma client after {max_retries} attempts") from last_exception

    def _initialize_collection(self) -> chromadb.Collection:
        """Initialize collection with proper schema and error handling"""
        max_retries = 3
        last_exception = None
        
        for attempt in range(max_retries):
            try:
                # Try with minimal metadata first
                metadata = {"hnsw:space": "cosine"} if attempt < 2 else None
                return self.client.get_or_create_collection(
                    name=COLLECTION_NAME,
                    embedding_function=self.embedding_fn,
                    metadata=metadata
                )
            except Exception as e:
                last_exception = e
                print(f"Collection initialization attempt {attempt + 1} failed: {e}")
                if attempt < max_retries - 1:
                    self._clean_chroma_db()
                    self.client = self._initialize_chroma_client()
        
        print("All collection initialization attempts failed")
        raise RuntimeError(f"Failed to initialize collection after {max_retries} attempts") from last_exception

    @staticmethod
    def extract_text(file_path: str) -> Optional[str]:
        """Extract text from various file formats with improved error handling"""
        path = Path(file_path)
        if not path.exists():
            print(f"File not found: {path}")
            return None

        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                with fitz.open(str(path)) as doc:
                    return "\n".join(page.get_text() for page in doc if page.get_text().strip())
            elif ext in (".docx", ".doc"):
                return "\n".join(p.text for p in docx.Document(str(path)).paragraphs if p.text.strip())
            elif ext == ".pptx":
                prs = Presentation(str(path))
                return "\n".join(
                    shape.text.strip() 
                    for slide in prs.slides 
                    for shape in slide.shapes 
                    if hasattr(shape, "text") and shape.text.strip()
                )
            elif ext == ".txt":
                return path.read_text(encoding='utf-8', errors='ignore')
            elif ext in (".csv", ".xls", ".xlsx"):
                df = pd.read_csv(path) if ext == ".csv" else pd.read_excel(path)
                return df.to_string(index=False)
            else:
                print(f"Unsupported file format: {ext}")
                return None
        except Exception as e:
            print(f"Error extracting text from {path.name}: {e}")
            return None

    def _generate_metadata(self, file_path: Path, chunk_idx: int) -> Dict[str, Any]:
        """Generate consistent metadata for chunks"""
        return {
            "source": file_path.name,
            "source_path": str(file_path),
            "chunk_idx": chunk_idx,
            "file_size": file_path.stat().st_size,
            "file_modified": file_path.stat().st_mtime
        }

    def ingest_file(self, file_path: str) -> int:
        """Ingest single file into ChromaDB with improved error handling"""
        file_path = Path(file_path)
        if not file_path.exists():
            print(f"File not found: {file_path}")
            return 0

        print(f"Processing {file_path.name}...")
        
        text = self.extract_text(str(file_path))
        if not text or not text.strip():
            print("No text extracted from file.")
            return 0

        try:
            chunks = self.splitter.split_text(text)
            if not chunks:
                print("No chunks produced from text.")
                return 0

            ids = [str(uuid.uuid4()) for _ in chunks]
            metadatas = [self._generate_metadata(file_path, i) for i in range(len(chunks))]

            self.collection.add(
                ids=ids,
                documents=chunks,
                metadatas=metadatas
            )
            print(f"Added {len(chunks)} chunks from {file_path.name}")
            return len(chunks)
        except Exception as e:
            print(f"Failed to ingest {file_path.name}: {e}")
            # Try to recover by reinitializing collection
            try:
                self.collection = self._initialize_collection()
                return 0
            except Exception as inner_e:
                print(f"Failed to recover after ingestion error: {inner_e}")
                return 0

    def ingest_all(self) -> int:
        """Ingest all valid files in uploads directory with progress tracking"""
        total = 0
        processed_files = 0
        
        for file in UPLOAD_DIR.iterdir():
            if file.is_file():
                processed_files += 1
                chunks_added = self.ingest_file(file)
                total += chunks_added
                print(f"Processed {processed_files}: {file.name} ({chunks_added} chunks)")
        
        print(f"\nIngestion complete. Processed {processed_files} files, total chunks added: {total}")
        return total

if __name__ == "__main__":
    try:
        print("Starting document ingestion service...")
        ingester = DocumentIngester()
        total_chunks = ingester.ingest_all()
        print(f"\nSuccessfully processed documents. Total chunks added: {total_chunks}")
    except Exception as e:
        print(f"\nFatal error in document ingestion: {e}")
        raise